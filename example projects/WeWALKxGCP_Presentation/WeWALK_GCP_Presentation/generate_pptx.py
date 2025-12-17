from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement

# Create a presentation
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Blank slide layout
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# Background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(248, 249, 250) # #F8F9FA

# Logos Header
# Centered logos
slide.shapes.add_picture('/Users/ozany/Desktop/WeWALK_GCP_Presentation/wewalk_logo_final.png', Inches(6.0), Inches(0.5), height=Inches(0.5))

# Separator X
textbox = slide.shapes.add_textbox(Inches(7.75), Inches(0.45), Inches(0.5), Inches(0.5))
tf = textbox.text_frame
p = tf.paragraphs[0]
p.text = "×"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(95, 99, 104) # #5f6368 (secondary color)
p.alignment = PP_ALIGN.CENTER

slide.shapes.add_picture('/Users/ozany/Desktop/WeWALK_GCP_Presentation/google_cloud_logo.png', Inches(8.25), Inches(0.4), height=Inches(0.7))

def set_shadow(shape):
    # Add a subtle shadow using OXML
    # This simulates box-shadow: 0 4px 24px rgba(0, 0, 0, 0.06);
    # spPr is a direct attribute on the element in newer python-pptx versions or accessible via xml
    # shape.element is a CT_Shape
    spPr = shape.element.spPr
    effectLst = spPr.get_or_add_effectLst()
    outerShdw = OxmlElement("a:outerShdw")
    outerShdw.set("blurRad", "240000") # 24pt roughly
    outerShdw.set("dist", "40000")     # 4pt roughly
    outerShdw.set("dir", "5400000")    # 90 degrees (down)
    outerShdw.set("algn", "ctr")       # center
    outerShdw.set("rotWithShape", "0")
    
    srgbClr = OxmlElement("a:srgbClr")
    srgbClr.set("val", "000000")
    alpha = OxmlElement("a:alpha")
    alpha.set("val", "6000") # 6% opacity
    srgbClr.append(alpha)
    outerShdw.append(srgbClr)
    effectLst.append(outerShdw)

# Helper function to create a card
def create_card(slide, x, y, width, height, title, icon_path, items, highlight=None, tech_tags=None, migration_note=None):
    # Card background (Rounded Rectangle)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = RGBColor(224, 224, 224) # #e0e0e0
    shape.line.width = Pt(1)
    
    # Set shadow
    set_shadow(shape)
    
    # Adjust rounded corners
    # The adjustment value is 0.0 to 1.0, representing the ratio of the shorter side.
    # We want approx 16px radius. 
    # 16px is roughly 0.16 inches. 
    # Shortest side is height (approx 3 inches). 0.16 / 3 = ~0.05
    shape.adjustments[0] = 0.05

    # Header Section
    header_y = y + 0.3
    
    # Icon
    if icon_path:
        slide.shapes.add_picture(icon_path, Inches(x + 0.3), Inches(header_y), width=Inches(0.4), height=Inches(0.4))

    # Title
    textbox = slide.shapes.add_textbox(Inches(x + 0.8), Inches(header_y), Inches(width - 1.0), Inches(0.5))
    tf = textbox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(26, 115, 232) # #1a73e8 (primary color)
    p.font.name = 'Arial'

    # Separator Line
    # line = slide.shapes.add_connector(
    #     MSO_CONNECTOR.STRAIGHT, 
    #     Inches(x), Inches(header_y + 0.5), Inches(x + width), Inches(header_y + 0.5)
    # )
    # line.line.color.rgb = RGBColor(232, 234, 237)
    # line.line.width = Pt(1)

    # Content Start Y
    current_y = header_y + 0.6

    # Highlight (if any)
    if highlight:
        textbox = slide.shapes.add_textbox(Inches(x + 0.3), Inches(current_y), Inches(width - 0.6), Inches(0.3))
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = highlight
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(66, 133, 244) # #4285f4 (accent color)
        p.font.bold = True
        p.font.name = 'Arial'
        current_y += 0.25

    # List items
    if items:
        textbox = slide.shapes.add_textbox(Inches(x + 0.3), Inches(current_y), Inches(width - 0.6), Inches(height - (current_y - y) - 0.2))
        tf = textbox.text_frame
        tf.word_wrap = True
        
        for item in items:
            p = tf.add_paragraph()
            # Handle bold text manually if possible, but python-pptx rich text is tricky in one go.
            # We'll just put the text.
            p.text = "• " + item.replace("<strong>", "").replace("</strong>", "")
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(95, 99, 104) # #5f6368
            p.font.name = 'Arial'
            p.space_after = Pt(6)
            p.line_spacing = 1.3 # Closer to 1.5 css line-height
            current_y += 0.4 # Approximate height increase

    # Tech Tags (for Backend card)
    if tech_tags:
        tag_x = x + 0.3
        tag_y = current_y + 0.1
        for tag in tech_tags:
            # Create a small rounded rect for the tag
            tag_width = Inches(0.8) # Estimate
            if len(tag) > 10: tag_width = Inches(1.2)
            
            ts = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(tag_x), Inches(tag_y), tag_width, Inches(0.25))
            ts.fill.solid()
            ts.fill.fore_color.rgb = RGBColor(232, 240, 254) # #e8f0fe
            ts.line.fill.background() # No border
            ts.adjustments[0] = 0.5 # Fully rounded ends
            
            tf = ts.text_frame
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            
            p = tf.paragraphs[0]
            p.text = tag
            p.font.size = Pt(8)
            p.font.color.rgb = RGBColor(26, 115, 232)
            p.font.bold = True
            p.font.name = 'Arial'
            p.alignment = PP_ALIGN.CENTER
            
            tag_x += tag_width + Inches(0.1)
        current_y += 0.4

    # Migration Note
    if migration_note:
        # Background for note
        note_height = Inches(0.6)
        note_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, 
            Inches(x + 0.3), Inches(current_y), Inches(width - 0.6), note_height
        )
        note_shape.fill.solid()
        note_shape.fill.fore_color.rgb = RGBColor(230, 244, 234) # #e6f4ea
        note_shape.line.fill.background()
        note_shape.adjustments[0] = 0.2
        
        # Text for note
        textbox = slide.shapes.add_textbox(Inches(x + 0.4), Inches(current_y), Inches(width - 0.8), note_height)
        tf = textbox.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = migration_note.replace("<strong>", "").replace("</strong>", "")
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(52, 168, 83) # #34a853 (success color)
        p.font.bold = True
        p.font.name = 'Arial'


# Grid Layout Calculation
# Slide width 16, height 9
# Margins: Left 1, Right 1, Top 1.5, Bottom 1
# 3 Columns, 2 Rows
# Gap: 0.4
# Card Width: (16 - 2 - 0.8) / 3 = 4.4
# Card Height: (9 - 2.5 - 0.4) / 2 = 3.05

col_width = 4.4
row_height = 3.2
gap = 0.4
start_x = 1.2
start_y = 1.8

# Row 1
# Card 1: Navigation
create_card(slide, start_x, start_y, col_width, row_height, 
            "Multimodal Navigation", 
            '/Users/ozany/Desktop/WeWALK_GCP_Presentation/icon_navigation.png',
            ["Step-by-step guidance using Google Maps & Places APIs.", 
             "Hands-free operation: Smart cane interaction keeps phone in pocket.",
             "Planned: Enhanced safety with sidewalk, crosswalk, and obstacle detection."],
            highlight="Tailored for Visually Impaired")

# Card 2: Voice Assistant
create_card(slide, start_x + col_width + gap, start_y, col_width, row_height, 
            "Smart Voice Assistant", 
            '/Users/ozany/Desktop/WeWALK_GCP_Presentation/icon_assistant.png',
            ["Currently utilizing DialogFlow for voice interactions.",
             "Moving to Gemini 2.5 Flash & Pro for next-gen conversational reasoning.",
             "Future: On-device models for offline voice experience on iOS/Android."],
            highlight="Transitioning to Vertex AI Agent Engine")

# Card 3: Computer Vision
create_card(slide, start_x + (col_width + gap) * 2, start_y, col_width, row_height, 
            "Situational Awareness", 
            '/Users/ozany/Desktop/WeWALK_GCP_Presentation/icon_ai.png',
            ["Object finding, door detection, and room layout understanding.",
             "Powered by Gemma and other vision models.",
             "Solves the \"last-mile\" navigation problem."],
            highlight="Advanced Computer Vision")

# Row 2
row2_y = start_y + row_height + gap

# Card 4: Voice Generation
create_card(slide, start_x, row2_y, col_width, row_height, 
            "Voice Generation", 
            '/Users/ozany/Desktop/WeWALK_GCP_Presentation/icon_ai.png', # Reusing AI icon as per HTML
            ["High-quality text-to-speech via Google TTS.",
             "Seamless online & offline operation."],
            highlight="Natural Speech Output")

# Card 5: Analytics
create_card(slide, start_x + col_width + gap, row2_y, col_width, row_height, 
            "Mobility Analytics", 
            '/Users/ozany/Desktop/WeWALK_GCP_Presentation/icon_analytics.png',
            ["Actively developing metrics for cane angle, swipes, and obstacle interactions.",
             "Exploring Google Cloud services for advanced data processing."],
            highlight="Under Development")

# Card 6: Backend (Infrastructure)
create_card(slide, start_x + (col_width + gap) * 2, row2_y, col_width, row_height, 
            "Backend & Operations", 
            '/Users/ozany/Desktop/WeWALK_GCP_Presentation/icon_cloudrun.png',
            [], # No list items, just tags and note
            tech_tags=["Cloud Run", "Cloud SQL", "Firestore", "Artifact Registry"],
            migration_note="Recently migrated from AWS to Google Cloud, leveraging GCP credits for scale and innovation.")

prs.save('/Users/ozany/Desktop/WeWALK_GCP_Presentation/presentation_editable.pptx')
