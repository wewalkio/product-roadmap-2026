from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def convert_to_dark_theme(pptx_path, output_path):
    prs = Presentation(pptx_path)

    # Dark Theme Colors
    BG_COLOR = RGBColor(32, 33, 36)      # #202124 (Dark Grey/Black)
    CARD_BG_COLOR = RGBColor(48, 49, 52) # #303134 (Lighter Dark Grey)
    TEXT_COLOR = RGBColor(232, 234, 237) # #e8eaed (Light Grey/White)
    PRIMARY_COLOR = RGBColor(138, 180, 248) # #8ab4f8 (Light Blue)
    ACCENT_COLOR = RGBColor(138, 180, 248)  # Same light blue for accents
    SECONDARY_TEXT = RGBColor(189, 193, 198) # #bdc1c6

    # Icon Mapping (Filename -> White version path)
    # We will try to match by position or order if possible, but replacing images in python-pptx 
    # without knowing the exact relationship is hard. 
    # However, we can iterate through shapes and if it's a picture, check its location/size 
    # to guess which icon it is.
    # Or, simpler: Just apply a "recolor" if possible? No.
    # We will try to replace based on approximate location if we can.
    
    icon_map = {
        # Approx X positions for Row 1
        (1.5, 2.0): '/Users/ozany/Desktop/WeWALK_GCP_Presentation/icon_navigation_white.png',
        (6.0, 2.0): '/Users/ozany/Desktop/WeWALK_GCP_Presentation/icon_assistant_white.png',
        (10.5, 2.0): '/Users/ozany/Desktop/WeWALK_GCP_Presentation/icon_ai_white.png',
        # Row 2
        (3.5, 5.5): '/Users/ozany/Desktop/WeWALK_GCP_Presentation/icon_ai_white.png', # Voice Gen
        (8.25, 5.5): '/Users/ozany/Desktop/WeWALK_GCP_Presentation/icon_analytics_white.png',
        (10.5, 5.5): '/Users/ozany/Desktop/WeWALK_GCP_Presentation/icon_cloudrun_white.png', # Backend
        # Header Logos
        (6.0, 0.5): '/Users/ozany/Desktop/WeWALK_GCP_Presentation/wewalk_logo_white.png',
        (8.25, 0.4): '/Users/ozany/Desktop/WeWALK_GCP_Presentation/google_cloud_logo_white.png'
    }

    for slide in prs.slides:
        # Set Slide Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

        for shape in slide.shapes:
            # 1. Text Boxes
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # Check current color to decide mapping? 
                        # Or just force all text to white/light blue?
                        # If it was blue (titles), make it Light Blue.
                        # If it was black/grey (body), make it White/Light Grey.
                        
                        # Heuristic: If font size > 12 (Titles), use Primary.
                        if run.font.size and run.font.size.pt > 12:
                             run.font.color.rgb = PRIMARY_COLOR
                        else:
                             run.font.color.rgb = TEXT_COLOR

            # 2. Shapes (Cards)
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                # Check if it's a rounded rectangle (likely a card)
                # We can check fill color if needed, but let's assume all rounded rects are cards
                if hasattr(shape, "fill"):
                    try:
                        # Change fill to dark card bg
                        shape.fill.solid() # Ensure solid
                        shape.fill.fore_color.rgb = CARD_BG_COLOR
                        
                        # Change line/border
                        shape.line.color.rgb = RGBColor(95, 99, 104) # #5f6368
                    except:
                        pass

            # 3. Pictures (Icons)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Try to replace with white version based on position
                # We check if the shape's left/top is close to our known coordinates
                # Tolerance of 0.5 inches
                x_in = shape.left.inches
                y_in = shape.top.inches
                
                for (tx, ty), new_img_path in icon_map.items():
                    if abs(x_in - tx) < 1.0 and abs(y_in - ty) < 1.0:
                        # Found a match!
                        # Replacing image in python-pptx is tricky. 
                        # The standard way is to add a new picture and remove the old one, 
                        # preserving z-order is hard.
                        # However, we can try to access the blip and change the rId? 
                        # Or just add new image on top and delete old?
                        
                        # Let's add new image at same pos/size
                        try:
                            new_pic = slide.shapes.add_picture(new_img_path, shape.left, shape.top, shape.width, shape.height)
                            # Move new pic to roughly correct z-order? 
                            # We can't easily reorder. But adding it last puts it on top, which is fine.
                            # Then delete the old one.
                            # Deleting shapes in python-pptx is not directly supported via .delete()
                            # We have to remove from element tree.
                            sp = shape.element
                            sp.getparent().remove(sp)
                        except Exception as e:
                            print(f"Failed to replace image at {x_in},{y_in}: {e}")
                        break

    prs.save(output_path)

if __name__ == "__main__":
    convert_to_dark_theme(
        '/Users/ozany/Desktop/WeWALKxGCP.pptx',
        '/Users/ozany/Desktop/WeWALKxGCP_Dark.pptx'
    )
